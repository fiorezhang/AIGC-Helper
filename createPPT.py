# -*- coding: UTF-8 -*-

import requests
import os
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt, Cm

from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import time
from bisect import bisect_left, bisect_right

#EMU to CM
img_unit_base = 360000

'''
ppt slides template
0 Title (presentation title slide)
1 Title and Content
2 Section Header (sometimes called Segue)
3 Two Content (side by side bullet textboxes)
4 Comparison (same but additional title for each side by side content box)
5 Title Only
6 Blank
7 Content with Caption
8 Picture with Caption
'''

def create_new_ppt_by_name(file_name_t):
    prs = Presentation(file_name_t)   # creat a new ppt
    return prs

def create_new_ppt():
    #prs = Presentation()   # creat a new ppt
    tmp_path = os.getcwd() + '\chatModels\\'
    prs = Presentation(pptx=os.path.join(tmp_path, 'default.pptx'))   # creat a new ppt
    return prs

def save_ppt(prs_r):
    file_path = os.getcwd() + '\ppt\\'
    curr_time = (str)((int)(time.time()))
    file_name = 'test_' + curr_time + '.pptx'
    file_full_name = file_path + file_name
    prs_r.save(file_full_name)
    return file_full_name

def check_one_slide_layout(slide):
    for shape in slide.placeholders:
        phf = shape.placeholder_format
        print(shape.placeholder_format)
        print(phf.idx)
        print(shape.name)
        print(phf.type)
        print(f"{phf.idx}--{shape.name}--{phf.type}")
        #shape.text = f"{phf.idx}--{shape.name}--{phf.type}"

def get_one_slide_layout(slide, idx_a, name_a):
    cc = 0
    for shape in slide.placeholders:
        phf = shape.placeholder_format
        idx_a.append(phf.idx)
        name_a.append(shape.name)
        cc+=1
    return cc

def get_slides(presentation):
    slides = presentation.slides
    slide_num = slides.len
    return slides, slide_num

# default is Title only page (5)
def add_one_slide(prs, SLD_LAYOUT_TITLE_AND_CONTENT = 5):
   slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
   slide = prs.slides.add_slide(slide_layout)
   return slide

def get_one_slide(slides_t, slide_index_t):
    return slides_t[slide_index_t]

def add_pic(slide, img, img_left, img_top, img_height):
    img_path = img
    left = Cm(img_left)
    top  = Cm(img_top)
    pic  = slide.shapes.add_picture(img_path, left, top, height=Cm(img_height))

def insert_pic(slide, index, img):
    placeholder = slide.placeholders[index]
    pic = placeholder.insert_picture(img)

def add_textbox(slide, index_t, text_t):
    tb_top = Cm(5.6)
    tb_left = Cm(1)
    tb_w = Cm(24)
    tb_h = Cm(10)
    
    text_box = slide.shapes.add_textbox(tb_left, tb_top, tb_w, tb_h)
    text_body = text_box.text_frame
   
    print('before filled\t', text_body.margin_bottom, text_body.margin_left, text_body.margin_right, text_body.margin_top, text_body.vertical_anchor)    

    text_body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_body.word_wrap = True
    text_body.text = text_t
    text_body.vertical_anchor = MSO_ANCHOR.MIDDLE

'''
text_rang_array=[[288, 324, 420, 484, 625, 729, 1056, 1260, 1600],
            [240, 285, 357, 396, 525, 621, 891, 1080, 1500]]
text_font_array=[[20, 19, 17, 16, 14, 13, 11, 10, 8],
           [20, 19, 17, 16, 14, 13, 11, 10, 8]]

'''
text_rang_array=[[248, 284, 280, 444, 585, 689, 1016, 1220, 1560],
            [220, 265, 337, 376, 505, 561, 871, 1060, 1480]]
text_font_array=[[20, 19, 17, 16, 14, 13, 11, 10, 8],
           [20, 19, 17, 16, 14, 13, 11, 10, 8]]

text_idx_flag = ')'
idx_max = 20
text_array = []

layout_title = 0
layout_slide = 3
t_title_sub = 'Edit by AIGC Helper'
t_text_body = 'Can use AIGC Helper generator pic and then put here?'

#input number of text and english or chinese
def text_fit(text_num_r, lang_r):
    print('text_fit:')
    print('text num is ', text_num_r, 'text lang is ', lang_r)
    lang_t = (int)(lang_r)
    if lang_t < 0 and lang_t > 1:    #English
        print('Text_fit error: english or Chinese?')
        return -1

    text_num_pos = bisect_left(text_rang_array[lang_t], text_num_r)
    if(text_num_pos < len(text_rang_array[lang_t])):
        text_font = text_font_array[lang_t][text_num_pos]
    print('pos = ', text_num_pos)
    print('font = ', text_font)
    return Pt(text_font)

def text_add(slide, index_r, text_r, text_len_r, lang_r, font_flag):
    print('text_add:')
    print('Input: len is', text_len_r, 'language is ', lang_r, 'font flag is ', font_flag)
    text_body = slide.placeholders[index_r].text_frame
    text_body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_body.word_wrap = True
    text_body.text = text_r
    if font_flag == 0:
        text_body.paragraphs[0].font.size = Pt(20)
    else:
        text_body.paragraphs[0].font.size = text_fit(text_len_r, lang_r)
    #text_body.fit_text(max_size=Pt(45))
    #text_body.fit_text()

def text_split(text_r, index_r, start_r):
    index_str = str(index_r)
    len_text = len(text_r)
    idx_t = index_r
    while(idx_t < idx_max):
        target_text = str(idx_t) + text_idx_flag
        len_b = text_r.find(target_text, start_r)
        idx_t += 1
        if len_b > 0:
            if (text_r[len_b-1] != '1') and (len_b > (start_r + 2)):  #remove 15 impact 5
                return len_b, idx_t
            else:
                continue
        elif len_b == 0:
            return 0, idx_t

    return len_text, idx_max

def title_split(text_r, start_r, end_r):
    target_text = text_r[start_r:end_r]
    len_text = len(target_text)
    len_b = target_text.find(',')
    return len_b

def write_title(ppt_r, title_text):
    slide = add_one_slide(ppt_h, layout_title)
    title = slide.shapes.title
    title.text = title_text
    subtitle = slide.placeholders[1]
    subtitle.text = 'Edit by AIGC Helper'

def write_slides(prs_r, text_str_r, text_flag_r):
    text_body = text_str_r[text_flag_r]
    text_body.lstrip( )
    text_array = text_body.split('\n')

    if text_flag_r == "Native":
        Lang_flag = 0
    elif text_flag_r == "Translated":
        Lang_flag = 1
    else:
        print("write_slides Error: ", text_flag_r)
        return

    slide = add_one_slide(prs_r, layout_title)
    title = slide.shapes.title
    title.text = text_array[0]
    subtitle = slide.placeholders[1]
    subtitle.text = 'Edit by AIGC Helper'

    array_idx = 2
    array_len = len(text_array)
    #print('array_len = ', array_len)
    while (array_idx < array_len):
        text_full_len = len(text_array[array_idx])
        #print('text_full_len = ', text_full_len)
        text_start = 0
        text_end = 0
        text_idx = 1
        while(text_end < text_full_len):
            text_end, text_idx = text_split(text_array[array_idx], text_idx, text_start)
            #print('text_start = ', text_start, 'text_end = ', text_end, 'text_idx = ', text_idx)
            #print(text_array[array_idx][text_start : text_end])
            if text_end!=0:
                curr_slide = add_one_slide(prs_r, layout_slide)
                text_add(curr_slide, 1, text_array[array_idx][text_start : text_end], text_end-text_start, Lang_flag, 1)
                text_add(curr_slide, 2, t_text_body, len(t_text_body), Lang_flag, 0)
                text_start = text_end
        array_idx += 1

if __name__ == '__main__':

    ppt_h = create_new_ppt()
    write_title(ppt_h,'this is a demo slides')

    write_slides(ppt_h, sample_y, 'Native')
    write_slides(ppt_h, sample_y, 'Translated')
    ppt_file_name = save_ppt(ppt_h)
    print('ppt file is ', ppt_file_name)
